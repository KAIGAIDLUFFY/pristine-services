from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create a presentation object
prs = Presentation()

# Define a function to create a slide with a title and content
def add_slide(prs, title, content):
    slide_layout = prs.slide_layouts[1]  # Use the layout with a title and content
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    content_placeholder = slide.placeholders[1]

    # Set the title
    title_placeholder.text = title

    # Set the content
    content_frame = content_placeholder.text_frame
    content_frame.text = content[0]
    for paragraph in content[1:]:
        p = content_frame.add_paragraph()
        p.text = paragraph
        p.level = 0

# Define the content for each slide
slides_content = [
    ("Pristine Cleaning Services Website", ["Project Documentation", "Author: Kaigai Brian Mwangi", "Date: July 2024"]),
    ("Introduction", ["Overview of the Project", "Purpose and Objectives"]),
    ("Background of the Study", ["Current Manual System", "Challenges Faced"]),
    ("Problem Statement", ["Issues with Manual System", "Need for Digital Transformation"]),
    ("Proposed Solution", ["Overview of the Web-Based System", "Key Features and Benefits"]),
    ("Objectives", ["Main Objective", "Specific Objectives:", "Online Booking and Management", "Data Storage and Retrieval", "Admin Functions", "Payment Integration"]),
    ("Justification", ["Advantages of the System", "Efficiency and Accuracy"]),
    ("Scope", ["System Capabilities", "Primary Users"]),
    ("Project Risks and Mitigation", ["Identified Risks", "Mitigation Strategies"]),
    ("Budget", ["Itemized Budget Overview", "Total Cost"]),
    ("Time Schedule", ["Task Breakdown and Timeline"]),
    ("Gantt Chart", ["Visual Representation of Timeline"]),
    ("System Design and Testing", ["Dashboard Designs", "Testing Procedures", "Implementation Strategy"]),
    ("Limitations and Recommendations", ["Project Limitations", "Recommendations for Future Work"]),
    ("Conclusion and Questions", ["Conclusion", "Open Floor for Questions"]),
]

# Add slides to the presentation
for title, content in slides_content:
    add_slide(prs, title, content)

# Apply a blue and white theme to the slides
for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text_frame = shape.text_frame
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(18)
                run.font.color.rgb = RGBColor(0, 0, 139)  # Dark blue color for text
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background

# Save the presentation with a unique name to avoid conflict
pptx_path = "/mnt/data/Pristine_Cleaning_Services_Presentation_v3.pptx"
prs.save(pptx_path)

pptx_path
